#!/usr/bin/env python3
"""macOS/Linux port of new_li_mathworks_pptx.ps1 (python-pptx, no PowerPoint needed).

Builds a Li's MathWorks Presentation PPTX from the same JSON slide spec as the
PowerShell generator: official template base, official layouts/placeholders,
Li preset geometry, ground-truth.md + speech.md artifacts.

Differences from the Windows generator (accepted, documented):
- No animations (python-pptx cannot author them). `builds` still land in notes.
- No shape grouping; Li pattern tagging uses shape names (`Li_<Pattern>_...`).
- No fill transparency (fills are the light palette already).

Usage:
  .venv/bin/python new_li_mathworks_pptx_mac.py SPEC.json OUT.pptx \
      [--template public|confidential] [--ground-truth PATH] [--speech PATH] \
      [--export-preview]

Preview export needs LibreOffice (`soffice`) and pypdfium2.
"""

import argparse
import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: F401  (PP_ALIGN kept for spec extensions)
from pptx.util import Pt

SKILL_ROOT = Path(__file__).resolve().parent.parent

PALETTE = {
    "Blue": RGBColor(0, 118, 168),
    "DeepBlue": RGBColor(0, 75, 135),
    "Cyan": RGBColor(0, 169, 224),
    "Orange": RGBColor(215, 136, 36),
    "Yellow": RGBColor(242, 169, 0),
    "Red": RGBColor(183, 48, 44),
    "Green": RGBColor(72, 162, 63),
    "Ink": RGBColor(31, 31, 31),
    "Muted": RGBColor(99, 101, 105),
    "LightBlue": RGBColor(236, 245, 248),
    "LightCyan": RGBColor(232, 247, 251),
    "LightYellow": RGBColor(255, 246, 218),
    "LightGreen": RGBColor(237, 247, 235),
    "LightRed": RGBColor(250, 234, 233),
    "LightGray": RGBColor(242, 244, 245),
}
LIGHT_FILLS = [PALETTE[k] for k in ("LightBlue", "LightYellow", "LightGreen", "LightCyan", "LightGray", "LightRed")]
ACCENT_LINES = [PALETTE[k] for k in ("Blue", "Orange", "Green", "Cyan", "Muted", "Red")]

# Content band: slides are 960x540 pt. The template title occupies the top band,
# and the footer/logo the bottom. Generated content must fill BAND_TOP..BAND_BOTTOM
# instead of clustering under the title and leaving the lower half empty.
SLIDE_W, SLIDE_H = 960.0, 540.0
BAND_TOP, BAND_BOTTOM = 132.0, 466.0
BAND_LEFT, BAND_RIGHT = 58.0, 902.0
BAND_H = BAND_BOTTOM - BAND_TOP
BAND_W = BAND_RIGHT - BAND_LEFT


def band_rows(count, gap, top=BAND_TOP, bottom=BAND_BOTTOM, max_height=None):
    """Split the content band into `count` rows.

    Rows share the band, but `max_height` stops a short list from being stretched
    into oversized frames with text floating inside them; whatever is left over is
    distributed so the block sits centred in the band rather than hugging the title.
    """
    count = max(1, count)
    height = (bottom - top - gap * (count - 1)) / count
    if max_height is not None and height > max_height:
        height = max_height
        used = height * count + gap * (count - 1)
        top = top + max(0.0, (bottom - top - used) / 2)
    return [(top + i * (height + gap), height) for i in range(count)]


def grid_shape(count):
    """Columns per row that avoid orphan cards (never 3+1)."""
    if count <= 3:
        return count
    if count == 4:
        return 2
    if count in (5, 6):
        return 3
    return 4


def typographic(text):
    """Punctuation a person typing in PowerPoint would get from autocorrect."""
    if not text:
        return text
    text = re.sub(r"(?<=\S) - (?=\S)", " — ", text)
    text = text.replace("'", "’")
    return text

BANNED_VISIBLE_TEXT = [
    r"(?i)use progressive reveal",
    r"(?i)speaker note",
    r"(?i)presenter note",
    r"(?i)build note",
    r"(?i)click to",
    r"(?i)animate this",
    r"(?i)visual check",
    r"(?i)\b(the workshop|the close|the deck|the slide)\s+should\b",
]
NON_VISIBLE_KEYS = {
    "reason", "talkTrack", "builds", "code", "coverTalkTrack", "sourceFacts",
    "missingInputs", "assetStatus", "placeholderStatus", "audience", "purpose",
}


GT_BEGIN = "<!-- generated:slide-ledger START (edits inside this block are overwritten) -->"
GT_END = "<!-- generated:slide-ledger END -->"


class SpecError(Exception):
    pass


# ---------- spec validation (port of Test-VisibleText / Test-Spec) ----------

def check_visible_text(node, path="spec"):
    if node is None:
        return
    if isinstance(node, str):
        for pattern in BANNED_VISIBLE_TEXT:
            if re.search(pattern, node):
                raise SpecError(f"Visible slide text appears to contain build/meta language at {path}: '{node}'")
        return
    if isinstance(node, (int, float, bool)):
        return
    if isinstance(node, list):
        for i, item in enumerate(node):
            check_visible_text(item, f"{path}[{i}]")
        return
    if isinstance(node, dict):
        for key, value in node.items():
            if key in NON_VISIBLE_KEYS:
                continue
            check_visible_text(value, f"{path}.{key}")


def locator_items(slide_spec, progress_labels):
    if slide_spec.get("locator"):
        return list(slide_spec["locator"])
    active = slide_spec.get("activeProgressLabel")
    if active and progress_labels:
        return [{"label": lab, "active": lab == active} for lab in progress_labels]
    return []


def resolve_asset(path_str, spec_dir):
    p = Path(path_str)
    if p.is_absolute():
        return p
    candidate = spec_dir / p
    if candidate.exists():
        return candidate
    return SKILL_ROOT / p


def validate_spec(spec, spec_dir):
    progress_labels = [str(x) for x in spec.get("progressLabels", [])]
    seen = set()
    for label in progress_labels:
        if label.lower() in seen:
            raise SpecError(f"Top-level progressLabels contains duplicate label '{label}'.")
        seen.add(label.lower())

    for i, slide_spec in enumerate(spec.get("slides", []), start=1):
        if not slide_spec.get("reason"):
            raise SpecError(f"Slide spec {i} is missing reason. Record why the slide exists before generating.")
        if not slide_spec.get("talkTrack"):
            raise SpecError(f"Slide spec {i} is missing talkTrack. Generate the speech/talk track before styling.")
        items = locator_items(slide_spec, progress_labels)
        if items:
            seen_loc, active_count = set(), 0
            for item in items:
                label = item if isinstance(item, str) else str(item.get("label", ""))
                if label.lower() in seen_loc:
                    raise SpecError(f"Slide spec {i} has duplicate progress/locator label '{label}'.")
                seen_loc.add(label.lower())
                if isinstance(item, dict) and item.get("active"):
                    active_count += 1
            if active_count != 1:
                raise SpecError(f"Slide spec {i} has {active_count} active progress/locator labels; expected exactly one.")
        active = slide_spec.get("activeProgressLabel")
        if active and progress_labels and active not in progress_labels:
            raise SpecError(f"Slide spec {i} activeProgressLabel '{active}' is not in top-level progressLabels.")
        if slide_spec.get("type") in ("image-evidence", "screenshot-evidence", "model-screenshot", "concept-artifact"):
            for image in slide_spec.get("images", []):
                path = image if isinstance(image, str) else str(image.get("path", ""))
                is_placeholder = isinstance(image, dict) and bool(image.get("placeholder"))
                if not is_placeholder and not resolve_asset(path, spec_dir).exists():
                    raise SpecError(f"Slide spec {i} image not found: {path}")
    return progress_labels


# ---------- shape helpers (ports of the PS helpers) ----------

def semantic_colors(role, index):
    r = (role or "").lower()
    if re.search(r"legacy|current", r):
        return PALETTE["LightGray"], PALETTE["Muted"]
    if re.search(r"model|workflow|architecture|scale", r):
        return PALETTE["LightBlue"], PALETTE["Blue"]
    if re.search(r"conversion|action|next|resolve", r):
        return PALETTE["LightYellow"], PALETTE["Orange"]
    if re.search(r"verification|validated|pass", r):
        return PALETTE["LightGreen"], PALETTE["Green"]
    if re.search(r"risk|block|blocked|blocker|gap", r):
        return PALETTE["LightRed"], PALETTE["Red"]
    if re.search(r"review|highlight|pending", r):
        return PALETTE["LightYellow"], PALETTE["Yellow"]
    return LIGHT_FILLS[index % len(LIGHT_FILLS)], ACCENT_LINES[index % len(ACCENT_LINES)]


def status_fill(text):
    """Colour a status cell by what it says.

    The vocabulary must cover the words a person actually writes about work that
    happened -- "fixed", "accepted", "delivered", "closed", "blocked" -- otherwise
    honest status words render colourless and authors reword them to chase a colour.
    """
    if not text:
        return None
    if re.search(r"(?i)\b(pass|passed|complete|completed|validated|verified|done|ok|ready|"
                 r"fixed|resolved|closed|accepted|approved|delivered|migrated|signed off)\b", text):
        return PALETTE["LightGreen"]
    if re.search(r"(?i)\b(risk|gap|fail|failed|blocker|blocked|missing|open|overdue|"
                 r"not started|rejected|stalled)\b", text):
        return PALETTE["LightRed"]
    if re.search(r"(?i)\b(review|pending|next|watch|in progress|partial|waiting|"
                 r"tbc|tbd|proposed|deferred|unconfirmed)\b", text):
        return PALETTE["LightYellow"]
    return None


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise SpecError(f"Template layout not found: {name}")


def _ph_type_name(shape):
    try:
        return shape.placeholder_format.type.name if shape.placeholder_format.type is not None else ""
    except Exception:
        return ""


def find_title(slide):
    for shape in slide.shapes:
        if _ph_type_name(shape) in ("TITLE", "CENTER_TITLE"):
            return shape
    for shape in slide.shapes:
        if shape.has_text_frame and "Title" in shape.name:
            return shape
    return None


def content_placeholders(slide):
    items = [s for s in slide.shapes if _ph_type_name(s) in ("BODY", "OBJECT")]
    return sorted(items, key=lambda s: s.left or 0)


def find_content(slide):
    phs = content_placeholders(slide)
    return phs[0] if phs else None


def tag(shape, pattern):
    try:
        shape.name = f"Li_{pattern}_{shape.shape_id}"
    except Exception:
        pass


def set_text(shape, text, size=0, bold=None, color=None):
    if shape is None:
        return
    tf = shape.text_frame
    tf.text = str(text)
    for para in tf.paragraphs:
        for run in para.runs:
            if size:
                run.font.size = Pt(size)
            if bold is not None:
                run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color


def add_textbox(slide, text, left, top, width, height, size, bold, color, allow_empty=False):
    if not allow_empty and not str(text or "").strip():
        return None
    shape = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(width), Pt(height))
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    lines = typographic(str(text)).split("\n")
    tf.text = lines[0]
    for extra in lines[1:]:
        tf.add_paragraph().text = extra
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return shape


def item_lines(items):
    lines = []
    for item in items or []:
        if isinstance(item, str):
            lines.append(item)
        else:
            label = str(item.get("label", ""))
            detail = str(item.get("detail", "") or "")
            lines.append(f"{label} - {detail}" if detail else label)
    return lines


def fill_band_with_text(shape, items, min_size=16.0, max_size=25.0):
    """Make a text placeholder occupy the content band instead of hugging the title.

    A short list rendered small at the top of the slide leaves the bottom half
    empty -- the single most common reason a generated deck looks unfinished.
    Fewer lines therefore means larger type, vertically centred in the band.
    """
    if shape is None:
        return 18.0
    count = max(1, len(item_lines(items)))
    # Scale type so a short list genuinely occupies the band. Middle-anchoring alone
    # just moves the empty space from below the text to above it.
    size = max(min_size, min(max_size, 118.0 / max(3, count)))
    try:
        # Resolve inherited geometry BEFORE writing any of it: setting only top/height
        # on an inherited placeholder materialises an <a:off>/<a:ext> with zeroed
        # left/width, which collapses the shape and loses the text entirely.
        left = shape.left if shape.left is not None else Pt(BAND_LEFT)
        width = shape.width if shape.width is not None else Pt(BAND_W)
        shape.left, shape.width = left, width
        shape.top = Pt(BAND_TOP)
        shape.height = Pt(BAND_H)
        # Top-anchored so the body stays attached to the headline it belongs to.
        shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        shape.text_frame.word_wrap = True
    except Exception:
        pass
    return size


def space_paragraphs(shape, points):
    for para in shape.text_frame.paragraphs[1:]:
        try:
            para.space_before = Pt(points)
        except Exception:
            pass


def set_body_paragraphs(shape, items, font_size=16):
    if shape is None:
        return
    tf = shape.text_frame
    lines = [typographic(line) for line in item_lines(items)]
    tf.text = lines[0] if lines else ""
    for line in lines[1:]:
        tf.add_paragraph().text = line
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = PALETTE["Ink"]


def add_card(slide, text, left, top, width, height, index, font_size=17, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(left), Pt(top), Pt(width), Pt(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or LIGHT_FILLS[index % len(LIGHT_FILLS)]
    shape.line.color.rgb = line or ACCENT_LINES[index % len(ACCENT_LINES)]
    shape.line.width = Pt(0.9)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Pt(12)
    tf.margin_top = tf.margin_bottom = Pt(8)
    tf.word_wrap = True
    # Top-anchored: PowerPoint centres autoshape text by default, which makes every
    # card in a row start at a different baseline whenever the text lengths differ.
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = typographic(str(text)).split("\n")
    tf.text = lines[0]
    for extra in lines[1:]:
        tf.add_paragraph().text = extra
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = PALETTE["Ink"]
    return shape


def two_line_card_fonts(card, first_size, second_size):
    paras = card.text_frame.paragraphs
    if len(paras) >= 2:
        for run in paras[0].runs:
            run.font.size = Pt(first_size)
        for run in paras[1].runs:
            run.font.size = Pt(second_size)
            run.font.bold = False
            run.font.color.rgb = PALETTE["Muted"]


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Pt(left), Pt(top), Pt(width), Pt(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALETTE["Orange"]
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_notes(slide, slide_spec):
    lines = []
    if slide_spec.get("reason"):
        lines.append(f"Reason: {slide_spec['reason']}")
    if slide_spec.get("talkTrack"):
        lines.append(f"Talk track: {slide_spec['talkTrack']}")
    if slide_spec.get("builds"):
        lines.append("Builds:")
        lines.extend(f"- {b}" for b in slide_spec["builds"])
    if lines:
        slide.notes_slide.notes_text_frame.text = "\n".join(lines)


def new_slide(prs, layout_name, slide_spec):
    slide = prs.slides.add_slide(get_layout(prs, layout_name))
    title = find_title(slide)
    if title is not None:
        title.text_frame.text = str(slide_spec.get("title", ""))
    return slide


def set_intro_in_content(slide, slide_spec, size=15):
    content = find_content(slide)
    if content is not None and slide_spec.get("intro"):
        set_text(content, slide_spec["intro"], size, False, PALETTE["Muted"])
        return content
    return content


# ---------- slide builders (geometry ported from the PS script) ----------

def add_agenda(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = find_content(slide)
    if content is None:
        content = add_textbox(slide, "", 48, 126, 848, 366, 18, False, PALETTE["Ink"], allow_empty=True)
    items = list(s.get("items", []))
    if s.get("intro"):
        items = [s["intro"]] + items
    size = fill_band_with_text(content, items, min_size=17.0, max_size=24.0)
    set_body_paragraphs(content, items, size)
    if content is not None:
        space_paragraphs(content, size * 0.7)


def add_content(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = find_content(slide)
    bullets = list(s.get("bullets", []))
    lines = item_lines(bullets)
    # A title plus two or three short bullets cannot fill the slide as a bullet list;
    # it only leaves a dead band. Promote it to stacked statement panels, which is
    # what a person does by hand when a slide looks empty.
    if 0 < len(lines) <= 3 and max((len(x) for x in lines), default=0) <= 120:
        if content is not None:
            try:
                content._element.getparent().remove(content._element)
            except Exception:
                pass
        longest = max((len(x) for x in lines), default=0)
        rows = band_rows(len(lines), 16, max_height=min(104.0, max(50.0, 34.0 + longest * 0.32)))
        for i, line in enumerate(lines):
            row_top, row_h = rows[i]
            card = add_card(slide, line, BAND_LEFT, row_top, BAND_W, row_h, i, 19,
                            fill=PALETTE["LightBlue"], line=PALETTE["Blue"])
            tag(card, "StatementPanel")
        return
    size = fill_band_with_text(content, bullets)
    set_body_paragraphs(content, bullets, size)
    if content is not None:
        space_paragraphs(content, size * 0.7)


def add_two_content(prs, s):
    slide = new_slide(prs, "Two Content", s)
    phs = content_placeholders(slide)
    left = phs[0] if len(phs) >= 1 else add_textbox(slide, "", 48, 126, 408, 366, 18, False, PALETTE["Ink"], allow_empty=True)
    right = phs[1] if len(phs) >= 2 else add_textbox(slide, "", 488, 126, 408, 366, 18, False, PALETTE["Ink"], allow_empty=True)
    set_body_paragraphs(left, [str(s.get("leftTitle", ""))] + list(s.get("leftBullets", [])), 16)
    set_body_paragraphs(right, [str(s.get("rightTitle", ""))] + list(s.get("rightBullets", [])), 16)
    for shape in (left, right):
        for run in shape.text_frame.paragraphs[0].runs:
            run.font.bold = True


def add_what_to_expect(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    outer = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Pt(BAND_LEFT - 4), Pt(BAND_TOP - 8),
        Pt(BAND_W + 8), Pt(BAND_H + 12))
    outer.fill.background()
    outer.line.color.rgb = PALETTE["Orange"]
    outer.line.width = Pt(0.9)
    outer.shadow.inherit = False
    tag(outer, "WhatToExpect")
    intro_h = 0.0
    if s.get("intro"):
        intro_h = 34.0
        box = add_textbox(slide, s["intro"], BAND_LEFT + 18, BAND_TOP + 8,
                          BAND_W - 36, intro_h, 15, False, PALETTE["Ink"])
        tag(box, "WhatToExpect")
    sections = s.get("sections") or s.get("items") or [
        {"label": "Best Practices", "detail": "Concepts and standards to reuse.", "role": "best-practice"},
        {"label": "Why should I do this?", "detail": "Motivation and risk context.", "role": "why"},
        {"label": "How do I do this?", "detail": "Implementation steps and examples.", "role": "how"},
        {"label": "Demo or Exercise", "detail": "Practice with a concrete artifact.", "role": "demo"},
        {"label": "Discussion", "detail": "Open questions and decisions.", "role": "discussion"},
    ]
    role_colors = {
        "best-practice": (RGBColor(244, 238, 249), RGBColor(112, 48, 160)),
        "extra-note": (RGBColor(248, 242, 252), RGBColor(180, 130, 218)),
        "why": (PALETTE["LightYellow"], PALETTE["Orange"]),
        "how": (PALETTE["LightGreen"], PALETTE["Green"]),
        "what": (PALETTE["LightBlue"], PALETTE["Blue"]),
        "demo": (PALETTE["LightCyan"], PALETTE["Cyan"]),
        "discussion": (PALETTE["LightRed"], PALETTE["Red"]),
    }
    # Bands share the framed area, so 5 or 8 sections both fit instead of overflowing.
    rows = band_rows(len(sections), 8, top=BAND_TOP + intro_h + 12, bottom=BAND_BOTTOM - 4, max_height=62)
    for i, section in enumerate(sections):
        label = section if isinstance(section, str) else str(section.get("label", ""))
        detail = "" if isinstance(section, str) else str(section.get("detail", "") or "")
        role = "" if isinstance(section, str) else str(section.get("role", "") or "")
        fill, line = role_colors.get(role, (LIGHT_FILLS[i % len(LIGHT_FILLS)], ACCENT_LINES[i % len(ACCENT_LINES)]))
        row_top, row_h = rows[i]
        text = f"{label}\n{detail}" if detail else label
        band = add_card(slide, text, BAND_LEFT + 18, row_top, BAND_W - 36, row_h, i, 13.5,
                        fill=fill, line=line)
        if detail:
            two_line_card_fonts(band, 13.5, 11)
        tag(band, "WhatToExpect")


def add_progress_rail(slide, items, active_label, left, top, width, row_height):
    for i, item in enumerate(items or []):
        label = item if isinstance(item, str) else str(item.get("label", ""))
        active = bool(item.get("active")) if isinstance(item, dict) else False
        if active_label:
            active = label == active_label
        fill = PALETTE["LightGreen"] if active else PALETTE["LightGray"]
        line = PALETTE["Green"] if active else PALETTE["Orange"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Pt(left), Pt(top + i * (row_height + 4)), Pt(width), Pt(row_height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5 if active else 0.8)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Pt(10)
        tf.margin_top = tf.margin_bottom = Pt(6)
        tf.word_wrap = True
        tf.text = label
        for run in tf.paragraphs[0].runs:
            run.font.name = "Arial"
            run.font.size = Pt(10.5)
            run.font.bold = active
            run.font.color.rgb = PALETTE["Ink"] if active else PALETTE["Muted"]
        tag(shape, "ProgressRail")


def add_progress_sidebar(prs, s, progress_labels):
    slide = new_slide(prs, "Title and Content", s)
    items = locator_items(s, progress_labels) or list(s.get("steps", []))
    rail_w = 226.0
    rail_rows = max(1, len(items))
    rail_h = min(34.0, (BAND_H - (rail_rows - 1) * 4) / rail_rows)
    add_progress_rail(slide, items, str(s.get("activeProgressLabel", "") or ""),
                      BAND_LEFT, BAND_TOP, rail_w, rail_h)
    main_left = BAND_LEFT + rail_w + 28
    main_w = BAND_RIGHT - main_left
    main_title = str(s.get("mainTitle") or s.get("activeProgressLabel") or "")
    header_h = 50.0 if main_title else 0.0
    if main_title:
        header = add_card(slide, main_title, main_left, BAND_TOP, main_w, header_h, 0, 20,
                          fill=PALETTE["LightGreen"], line=PALETTE["Green"])
        header.line.width = Pt(1.5)
        tag(header, "ProgressSidebar")
    cards = s.get("cards") or s.get("items") or ([{"label": s["intro"], "role": "highlight"}] if s.get("intro") else [])
    if not cards:
        return
    per_row = 1 if len(cards) <= 2 else 2
    row_count = -(-len(cards) // per_row)
    rows = band_rows(row_count, 14, top=BAND_TOP + header_h + (14 if main_title else 0), max_height=132)
    gap = 18
    for i, card in enumerate(cards):
        label = card if isinstance(card, str) else str(card.get("label", ""))
        detail = "" if isinstance(card, str) else str(card.get("detail", "") or "")
        role = "" if isinstance(card, str) else str(card.get("role", "") or "")
        fill, line = semantic_colors(role, i)
        text = f"{label}\n{detail}" if detail else label
        row, col = divmod(i, per_row)
        row_top, row_h = rows[row]
        in_row = min(per_row, len(cards) - row * per_row)
        cw = (main_w - gap * (in_row - 1)) / in_row
        shape = add_card(slide, text, main_left + col * (cw + gap), row_top, cw, row_h, i, 15,
                         fill=fill, line=line)
        if detail:
            two_line_card_fonts(shape, 15, 12)
        tag(shape, "ProgressSidebar")


def add_image_evidence(prs, s, pattern, spec_dir):
    """Insert real pictures. A missing asset is named in the body text, never drawn
    as an empty rectangle -- an empty box reads as a broken slide."""
    slide = new_slide(prs, "Title and Content", s)
    body = find_content(slide)
    images = list(s.get("images", []))
    resolved = []
    missing = []
    for image in images:
        path = image if isinstance(image, str) else str(image.get("path", ""))
        caption = "" if isinstance(image, str) else str(image.get("caption", "") or "")
        target = resolve_asset(path, spec_dir) if path else None
        if target is not None and target.exists():
            resolved.append((target, caption))
        else:
            missing.append(caption or "screenshot or evidence asset")

    lines = []
    if s.get("intro"):
        lines.append(str(s["intro"]))
    for label, detail in outline_items(s):
        lines.append(f"{label} — {detail}" if label and detail else (label or detail))
    for item in missing:
        lines.append(f"To add: {item}")

    if resolved:
        # Picture occupies the body area; supporting text goes above it if short.
        if body is not None:
            try:
                left, top = body.left, body.top
                width, height = body.width, body.height
                body._element.getparent().remove(body._element)
            except Exception:
                left, top, width, height = Pt(BAND_LEFT), Pt(BAND_TOP), Pt(BAND_W), Pt(BAND_H)
        else:
            left, top, width, height = Pt(BAND_LEFT), Pt(BAND_TOP), Pt(BAND_W), Pt(BAND_H)
        count = len(resolved)
        gap = Pt(18)
        each = int((int(width) - int(gap) * (count - 1)) / count)
        for i, (target, caption) in enumerate(resolved):
            pic = slide.shapes.add_picture(str(target), int(left) + i * (each + int(gap)),
                                           int(top), each, int(height))
            pic.line.color.rgb = PALETTE["Muted"]
            pic.line.width = Pt(0.75)
            if caption:
                lines.append(caption)
        if lines:
            box = add_textbox(slide, "\n".join(lines), BAND_LEFT, BAND_BOTTOM + 4,
                              BAND_W, 26, 11, False, PALETTE["Muted"])
            tag(box, pattern)
        return slide

    if body is not None and lines:
        tf = body.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = typographic(lines[0])
        for line in lines[1:]:
            tf.add_paragraph().text = typographic(line)
    return slide


def status_fill(text):
    """Colour a status cell by what it says.

    The vocabulary must cover the words a person actually writes about work that
    happened -- "fixed", "accepted", "delivered", "closed", "blocked" -- otherwise
    honest status words render colourless and authors reword them to chase a colour.
    """
    if not text:
        return None
    if re.search(r"(?i)\b(pass|passed|complete|completed|validated|verified|done|ok|ready|"
                 r"fixed|resolved|closed|accepted|approved|delivered|migrated|signed off)\b", text):
        return PALETTE["LightGreen"]
    if re.search(r"(?i)\b(risk|gap|fail|failed|blocker|blocked|missing|open|overdue|"
                 r"not started|rejected|stalled)\b", text):
        return PALETTE["LightRed"]
    if re.search(r"(?i)\b(review|pending|next|watch|in progress|partial|waiting|"
                 r"tbc|tbd|proposed|deferred|unconfirmed)\b", text):
        return PALETTE["LightYellow"]
    return None


def get_layout(prs, name):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == name:
                return layout
    raise SpecError(f"Template layout not found: {name}")


def _ph_type_name(shape):
    try:
        return shape.placeholder_format.type.name if shape.placeholder_format.type is not None else ""
    except Exception:
        return ""


def find_title(slide):
    for shape in slide.shapes:
        if _ph_type_name(shape) in ("TITLE", "CENTER_TITLE"):
            return shape
    for shape in slide.shapes:
        if shape.has_text_frame and "Title" in shape.name:
            return shape
    return None


def content_placeholders(slide):
    items = [s for s in slide.shapes if _ph_type_name(s) in ("BODY", "OBJECT")]
    return sorted(items, key=lambda s: s.left or 0)


def find_content(slide):
    phs = content_placeholders(slide)
    return phs[0] if phs else None


def tag(shape, pattern):
    try:
        shape.name = f"Li_{pattern}_{shape.shape_id}"
    except Exception:
        pass


def set_text(shape, text, size=0, bold=None, color=None):
    if shape is None:
        return
    tf = shape.text_frame
    tf.text = str(text)
    for para in tf.paragraphs:
        for run in para.runs:
            if size:
                run.font.size = Pt(size)
            if bold is not None:
                run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color


def add_textbox(slide, text, left, top, width, height, size, bold, color, allow_empty=False):
    if not allow_empty and not str(text or "").strip():
        return None
    shape = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(width), Pt(height))
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    lines = typographic(str(text)).split("\n")
    tf.text = lines[0]
    for extra in lines[1:]:
        tf.add_paragraph().text = extra
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
    return shape


def item_lines(items):
    lines = []
    for item in items or []:
        if isinstance(item, str):
            lines.append(item)
        else:
            label = str(item.get("label", ""))
            detail = str(item.get("detail", "") or "")
            lines.append(f"{label} - {detail}" if detail else label)
    return lines


def fill_band_with_text(shape, items, min_size=16.0, max_size=25.0):
    """Make a text placeholder occupy the content band instead of hugging the title.

    A short list rendered small at the top of the slide leaves the bottom half
    empty -- the single most common reason a generated deck looks unfinished.
    Fewer lines therefore means larger type, vertically centred in the band.
    """
    if shape is None:
        return 18.0
    count = max(1, len(item_lines(items)))
    # Scale type so a short list genuinely occupies the band. Middle-anchoring alone
    # just moves the empty space from below the text to above it.
    size = max(min_size, min(max_size, 118.0 / max(3, count)))
    try:
        # Resolve inherited geometry BEFORE writing any of it: setting only top/height
        # on an inherited placeholder materialises an <a:off>/<a:ext> with zeroed
        # left/width, which collapses the shape and loses the text entirely.
        left = shape.left if shape.left is not None else Pt(BAND_LEFT)
        width = shape.width if shape.width is not None else Pt(BAND_W)
        shape.left, shape.width = left, width
        shape.top = Pt(BAND_TOP)
        shape.height = Pt(BAND_H)
        # Top-anchored so the body stays attached to the headline it belongs to.
        shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
        shape.text_frame.word_wrap = True
    except Exception:
        pass
    return size


def space_paragraphs(shape, points):
    for para in shape.text_frame.paragraphs[1:]:
        try:
            para.space_before = Pt(points)
        except Exception:
            pass


def set_body_paragraphs(shape, items, font_size=16):
    if shape is None:
        return
    tf = shape.text_frame
    lines = [typographic(line) for line in item_lines(items)]
    tf.text = lines[0] if lines else ""
    for line in lines[1:]:
        tf.add_paragraph().text = line
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = PALETTE["Ink"]


def add_card(slide, text, left, top, width, height, index, font_size=17, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(left), Pt(top), Pt(width), Pt(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or LIGHT_FILLS[index % len(LIGHT_FILLS)]
    shape.line.color.rgb = line or ACCENT_LINES[index % len(ACCENT_LINES)]
    shape.line.width = Pt(0.9)
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Pt(12)
    tf.margin_top = tf.margin_bottom = Pt(8)
    tf.word_wrap = True
    # Top-anchored: PowerPoint centres autoshape text by default, which makes every
    # card in a row start at a different baseline whenever the text lengths differ.
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = typographic(str(text)).split("\n")
    tf.text = lines[0]
    for extra in lines[1:]:
        tf.add_paragraph().text = extra
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = PALETTE["Ink"]
    return shape


def two_line_card_fonts(card, first_size, second_size):
    paras = card.text_frame.paragraphs
    if len(paras) >= 2:
        for run in paras[0].runs:
            run.font.size = Pt(first_size)
        for run in paras[1].runs:
            run.font.size = Pt(second_size)
            run.font.bold = False
            run.font.color.rgb = PALETTE["Muted"]


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Pt(left), Pt(top), Pt(width), Pt(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALETTE["Orange"]
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_notes(slide, slide_spec):
    lines = []
    if slide_spec.get("reason"):
        lines.append(f"Reason: {slide_spec['reason']}")
    if slide_spec.get("talkTrack"):
        lines.append(f"Talk track: {slide_spec['talkTrack']}")
    if slide_spec.get("builds"):
        lines.append("Builds:")
        lines.extend(f"- {b}" for b in slide_spec["builds"])
    if lines:
        slide.notes_slide.notes_text_frame.text = "\n".join(lines)


def new_slide(prs, layout_name, slide_spec):
    slide = prs.slides.add_slide(get_layout(prs, layout_name))
    title = find_title(slide)
    if title is not None:
        title.text_frame.text = str(slide_spec.get("title", ""))
    return slide


def set_intro_in_content(slide, slide_spec, size=15):
    content = find_content(slide)
    if content is not None and slide_spec.get("intro"):
        set_text(content, slide_spec["intro"], size, False, PALETTE["Muted"])
        return content
    return content


# ---------- slide builders (geometry ported from the PS script) ----------

def add_agenda(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = find_content(slide)
    if content is None:
        content = add_textbox(slide, "", 48, 126, 848, 366, 18, False, PALETTE["Ink"], allow_empty=True)
    items = list(s.get("items", []))
    if s.get("intro"):
        items = [s["intro"]] + items
    size = fill_band_with_text(content, items, min_size=17.0, max_size=24.0)
    set_body_paragraphs(content, items, size)
    if content is not None:
        space_paragraphs(content, size * 0.7)


def add_content(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = find_content(slide)
    bullets = list(s.get("bullets", []))
    lines = item_lines(bullets)
    # A title plus two or three short bullets cannot fill the slide as a bullet list;
    # it only leaves a dead band. Promote it to stacked statement panels, which is
    # what a person does by hand when a slide looks empty.
    if 0 < len(lines) <= 3 and max((len(x) for x in lines), default=0) <= 120:
        if content is not None:
            try:
                content._element.getparent().remove(content._element)
            except Exception:
                pass
        longest = max((len(x) for x in lines), default=0)
        rows = band_rows(len(lines), 16, max_height=min(104.0, max(50.0, 34.0 + longest * 0.32)))
        for i, line in enumerate(lines):
            row_top, row_h = rows[i]
            card = add_card(slide, line, BAND_LEFT, row_top, BAND_W, row_h, i, 19,
                            fill=PALETTE["LightBlue"], line=PALETTE["Blue"])
            tag(card, "StatementPanel")
        return
    size = fill_band_with_text(content, bullets)
    set_body_paragraphs(content, bullets, size)
    if content is not None:
        space_paragraphs(content, size * 0.7)


def add_two_content(prs, s):
    slide = new_slide(prs, "Two Content", s)
    phs = content_placeholders(slide)
    left = phs[0] if len(phs) >= 1 else add_textbox(slide, "", 48, 126, 408, 366, 18, False, PALETTE["Ink"], allow_empty=True)
    right = phs[1] if len(phs) >= 2 else add_textbox(slide, "", 488, 126, 408, 366, 18, False, PALETTE["Ink"], allow_empty=True)
    set_body_paragraphs(left, [str(s.get("leftTitle", ""))] + list(s.get("leftBullets", [])), 16)
    set_body_paragraphs(right, [str(s.get("rightTitle", ""))] + list(s.get("rightBullets", [])), 16)
    for shape in (left, right):
        for run in shape.text_frame.paragraphs[0].runs:
            run.font.bold = True


def add_what_to_expect(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    outer = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Pt(BAND_LEFT - 4), Pt(BAND_TOP - 8),
        Pt(BAND_W + 8), Pt(BAND_H + 12))
    outer.fill.background()
    outer.line.color.rgb = PALETTE["Orange"]
    outer.line.width = Pt(0.9)
    outer.shadow.inherit = False
    tag(outer, "WhatToExpect")
    intro_h = 0.0
    if s.get("intro"):
        intro_h = 34.0
        box = add_textbox(slide, s["intro"], BAND_LEFT + 18, BAND_TOP + 8,
                          BAND_W - 36, intro_h, 15, False, PALETTE["Ink"])
        tag(box, "WhatToExpect")
    sections = s.get("sections") or s.get("items") or [
        {"label": "Best Practices", "detail": "Concepts and standards to reuse.", "role": "best-practice"},
        {"label": "Why should I do this?", "detail": "Motivation and risk context.", "role": "why"},
        {"label": "How do I do this?", "detail": "Implementation steps and examples.", "role": "how"},
        {"label": "Demo or Exercise", "detail": "Practice with a concrete artifact.", "role": "demo"},
        {"label": "Discussion", "detail": "Open questions and decisions.", "role": "discussion"},
    ]
    role_colors = {
        "best-practice": (RGBColor(244, 238, 249), RGBColor(112, 48, 160)),
        "extra-note": (RGBColor(248, 242, 252), RGBColor(180, 130, 218)),
        "why": (PALETTE["LightYellow"], PALETTE["Orange"]),
        "how": (PALETTE["LightGreen"], PALETTE["Green"]),
        "what": (PALETTE["LightBlue"], PALETTE["Blue"]),
        "demo": (PALETTE["LightCyan"], PALETTE["Cyan"]),
        "discussion": (PALETTE["LightRed"], PALETTE["Red"]),
    }
    # Bands share the framed area, so 5 or 8 sections both fit instead of overflowing.
    rows = band_rows(len(sections), 8, top=BAND_TOP + intro_h + 12, bottom=BAND_BOTTOM - 4, max_height=62)
    for i, section in enumerate(sections):
        label = section if isinstance(section, str) else str(section.get("label", ""))
        detail = "" if isinstance(section, str) else str(section.get("detail", "") or "")
        role = "" if isinstance(section, str) else str(section.get("role", "") or "")
        fill, line = role_colors.get(role, (LIGHT_FILLS[i % len(LIGHT_FILLS)], ACCENT_LINES[i % len(ACCENT_LINES)]))
        row_top, row_h = rows[i]
        text = f"{label}\n{detail}" if detail else label
        band = add_card(slide, text, BAND_LEFT + 18, row_top, BAND_W - 36, row_h, i, 13.5,
                        fill=fill, line=line)
        if detail:
            two_line_card_fonts(band, 13.5, 11)
        tag(band, "WhatToExpect")


def add_progress_rail(slide, items, active_label, left, top, width, row_height):
    for i, item in enumerate(items or []):
        label = item if isinstance(item, str) else str(item.get("label", ""))
        active = bool(item.get("active")) if isinstance(item, dict) else False
        if active_label:
            active = label == active_label
        fill = PALETTE["LightGreen"] if active else PALETTE["LightGray"]
        line = PALETTE["Green"] if active else PALETTE["Orange"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Pt(left), Pt(top + i * (row_height + 4)), Pt(width), Pt(row_height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5 if active else 0.8)
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Pt(10)
        tf.margin_top = tf.margin_bottom = Pt(6)
        tf.word_wrap = True
        tf.text = label
        for run in tf.paragraphs[0].runs:
            run.font.name = "Arial"
            run.font.size = Pt(10.5)
            run.font.bold = active
            run.font.color.rgb = PALETTE["Ink"] if active else PALETTE["Muted"]
        tag(shape, "ProgressRail")


def add_progress_sidebar(prs, s, progress_labels):
    slide = new_slide(prs, "Title and Content", s)
    items = locator_items(s, progress_labels) or list(s.get("steps", []))
    rail_w = 226.0
    rail_rows = max(1, len(items))
    rail_h = min(34.0, (BAND_H - (rail_rows - 1) * 4) / rail_rows)
    add_progress_rail(slide, items, str(s.get("activeProgressLabel", "") or ""),
                      BAND_LEFT, BAND_TOP, rail_w, rail_h)
    main_left = BAND_LEFT + rail_w + 28
    main_w = BAND_RIGHT - main_left
    main_title = str(s.get("mainTitle") or s.get("activeProgressLabel") or "")
    header_h = 50.0 if main_title else 0.0
    if main_title:
        header = add_card(slide, main_title, main_left, BAND_TOP, main_w, header_h, 0, 20,
                          fill=PALETTE["LightGreen"], line=PALETTE["Green"])
        header.line.width = Pt(1.5)
        tag(header, "ProgressSidebar")
    cards = s.get("cards") or s.get("items") or ([{"label": s["intro"], "role": "highlight"}] if s.get("intro") else [])
    if not cards:
        return
    per_row = 1 if len(cards) <= 2 else 2
    row_count = -(-len(cards) // per_row)
    rows = band_rows(row_count, 14, top=BAND_TOP + header_h + (14 if main_title else 0), max_height=132)
    gap = 18
    for i, card in enumerate(cards):
        label = card if isinstance(card, str) else str(card.get("label", ""))
        detail = "" if isinstance(card, str) else str(card.get("detail", "") or "")
        role = "" if isinstance(card, str) else str(card.get("role", "") or "")
        fill, line = semantic_colors(role, i)
        text = f"{label}\n{detail}" if detail else label
        row, col = divmod(i, per_row)
        row_top, row_h = rows[row]
        in_row = min(per_row, len(cards) - row * per_row)
        cw = (main_w - gap * (in_row - 1)) / in_row
        shape = add_card(slide, text, main_left + col * (cw + gap), row_top, cw, row_h, i, 15,
                         fill=fill, line=line)
        if detail:
            two_line_card_fonts(shape, 15, 12)
        tag(shape, "ProgressSidebar")


def add_image_evidence(prs, s, pattern, spec_dir):
    slide = new_slide(prs, "Title and Content", s)
    if s.get("intro"):
        box = add_textbox(slide, s["intro"], 78, 135, 780, 34, 15, False, PALETTE["Muted"])
        tag(box, pattern)
    images = list(s.get("images", []))
    callouts = list(s.get("callouts", []))
    top = BAND_TOP + (30 if s.get("intro") else 0)
    # Media takes the full band width when there is nothing beside it; otherwise a
    # 62/38 split, and the callout rail stretches to the media height so the right
    # side is never left stranded.
    rail_w = 250.0 if callouts else 0.0
    media_w = BAND_W - (rail_w + 24 if callouts else 0)
    caption_h = 22.0
    media_h = BAND_BOTTOM - top - (caption_h if any(
        isinstance(im, dict) and im.get("caption") for im in images) else 0)
    has_real_image = any(
        (resolve_asset(im if isinstance(im, str) else str(im.get("path", "")), spec_dir).exists())
        for im in images if (im if isinstance(im, str) else im.get("path")))
    if not has_real_image:
        media_h = min(media_h, 158.0)

    count = max(1, len(images))
    per_media_w = (media_w - 18 * (count - 1)) / count
    for i, image in enumerate(images):
        path = image if isinstance(image, str) else str(image.get("path", ""))
        caption = "" if isinstance(image, str) else str(image.get("caption", "") or "")
        resolved = resolve_asset(path, spec_dir) if path else None
        left = BAND_LEFT + i * (per_media_w + 18)
        if resolved is not None and resolved.exists():
            pic = slide.shapes.add_picture(str(resolved), Pt(left), Pt(top), Pt(per_media_w), Pt(media_h))
            pic.line.color.rgb = PALETTE["Muted"]
            pic.line.width = Pt(0.75)
            tag(pic, pattern)
        else:
            # An unavailable figure is a hairline frame with ONE label -- never the
            # word "Placeholder" plus the same sentence repeated as a caption.
            frame = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Pt(left), Pt(top), Pt(per_media_w), Pt(media_h))
            frame.fill.background()
            frame.line.color.rgb = PALETTE["Muted"]
            frame.line.width = Pt(0.75)
            frame.shadow.inherit = False
            tag(frame, pattern)
            label = caption or "Screenshot or evidence asset to be supplied"
            box = add_textbox(slide, f"[ to add ]  {label}", left + 16, top + 16,
                              per_media_w - 32, 40, 12.5, False, PALETTE["Muted"])
            tag(box, pattern)
            continue
        if caption:
            box = add_textbox(slide, caption, left, top + media_h + 4, per_media_w,
                              caption_h, 11, False, PALETTE["Muted"])
            tag(box, pattern)

    if callouts:
        rail_left = BAND_RIGHT - rail_w
        rows = band_rows(len(callouts), 14, top=top, bottom=top + media_h, max_height=92)
        for i, callout in enumerate(callouts):
            label = callout if isinstance(callout, str) else str(callout.get("label", ""))
            role = "" if isinstance(callout, str) else str(callout.get("role", "") or "")
            fill, line = semantic_colors(role, i)
            row_top, row_h = rows[i]
            box = add_card(slide, label, rail_left, row_top, rail_w, row_h, i, 13, fill=fill, line=line)
            tag(box, pattern)


def add_comparison_table(prs, s, pattern):
    slide = new_slide(prs, "Title and Content", s)
    if s.get("intro"):
        box = add_textbox(slide, s["intro"], 78, 135, 790, 34, 15, False, PALETTE["Muted"])
        tag(box, pattern)
    headers = list(s.get("headers", [])) or ["Item", "Evidence", "Status", "Next"]
    rows = list(s.get("rows", [])) or list(s.get("items", []))
    row_count = max(2, len(rows) + 1)
    col_count = max(2, len(headers))
    band_top = BAND_TOP + (30 if s.get("intro") else 0)
    band_bottom = BAND_BOTTOM - (44 if s.get("callout") else 0)
    header_h = 30.0
    n_rows = max(1, len(rows))
    # Rows share the band but never balloon: a 2-row table stretched over the whole
    # band looks as wrong as one crammed at the top. Cap, then centre what's left.
    body_h = min(64.0, max(30.0, (band_bottom - band_top - header_h) / n_rows))
    total_h = header_h + body_h * n_rows
    table_top = band_top + max(0.0, (band_bottom - band_top - total_h) / 2)
    table_shape = slide.shapes.add_table(
        row_count, col_count, Pt(BAND_LEFT), Pt(table_top), Pt(BAND_W), Pt(total_h))
    tag(table_shape, pattern)
    table = table_shape.table
    # Header height follows its single line of text; data rows share the remaining band.
    table.rows[0].height = Pt(header_h)
    for r in range(1, row_count):
        table.rows[r].height = Pt(body_h)
    for c in range(col_count):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PALETTE["LightBlue"]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = cell.margin_bottom = Pt(2)
        cell.text_frame.text = str(headers[c]) if c < len(headers) else ""
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = PALETTE["DeepBlue"]
    for r, row in enumerate(rows):
        if isinstance(row, str):
            cells = [part.strip() for part in row.split("|")]
        elif isinstance(row, dict) and row.get("cells"):
            cells = [str(v) for v in row["cells"]]
        elif isinstance(row, dict):
            cells = [str(row.get(str(h), row.get(str(h).lower(), ""))) for h in headers]
        else:
            cells = [str(v) for v in row]
        for c in range(col_count):
            value = cells[c] if c < len(cells) else ""
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.text = typographic(value)
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(11)
                run.font.color.rgb = PALETTE["Ink"]
            header_text = str(headers[c]) if c < len(headers) else ""
            is_status_col = bool(re.search(r"(?i)status|result|outcome", header_text))
            fill = status_fill(value)
            if is_status_col and fill is not None:
                cell.fill.fore_color.rgb = fill
            elif r % 2 == 1:
                cell.fill.fore_color.rgb = PALETTE["LightGray"]
    if s.get("callout"):
        callout = add_card(slide, str(s["callout"]), BAND_LEFT + 60, BAND_BOTTOM - 36,
                           BAND_W - 120, 34, 2, 14,
                           fill=PALETTE["LightYellow"], line=PALETTE["Orange"])
        tag(callout, pattern)


def add_concept_artifact(prs, s, spec_dir):
    slide = new_slide(prs, "Title and Content", s)
    content = set_intro_in_content(slide, s)
    if content is not None:
        tag(content, "ConceptArtifact")
    concepts = s.get("concepts") or s.get("items") or []
    top = BAND_TOP + (30 if s.get("intro") else 0)
    left_w = 372.0
    callouts = list(s.get("callouts", []))
    callout_band = 96.0 if callouts else 0.0
    art_bottom = BAND_BOTTOM - callout_band
    if concepts:
        rows = band_rows(len(concepts), 14, top=top, bottom=BAND_BOTTOM, max_height=96)
        for i, item in enumerate(concepts):
            label = item if isinstance(item, str) else str(item.get("label", ""))
            detail = "" if isinstance(item, str) else str(item.get("detail", "") or "")
            role = "model" if isinstance(item, str) else str(item.get("role", "") or "")
            fill, line = semantic_colors(role, i)
            text = f"{label}\n{detail}" if detail else label
            row_top, row_h = rows[i]
            card = add_card(slide, text, BAND_LEFT, row_top, left_w, row_h, i, 15, fill=fill, line=line)
            if detail:
                two_line_card_fonts(card, 15, 12)
            tag(card, "ConceptArtifact")
    art_left = BAND_LEFT + left_w + 24
    art_w = BAND_RIGHT - art_left
    artifact = None
    images = s.get("images", [])
    if images:
        image = images[0]
        path = image if isinstance(image, str) else str(image.get("path", ""))
        resolved = resolve_asset(path, spec_dir) if path else None
        if resolved is not None and resolved.exists():
            artifact = slide.shapes.add_picture(
                str(resolved), Pt(art_left), Pt(top), Pt(art_w), Pt(art_bottom - top))
            artifact.line.color.rgb = PALETTE["Muted"]
            artifact.line.width = Pt(0.75)
    if artifact is None:
        label = str(s.get("artifactLabel", "Artifact or evidence"))
        artifact = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Pt(art_left), Pt(top), Pt(art_w), Pt(art_bottom - top))
        artifact.fill.background()
        artifact.line.color.rgb = PALETTE["Muted"]
        artifact.line.width = Pt(0.75)
        artifact.shadow.inherit = False
        note = add_textbox(slide, f"[ to add ]  {label}", art_left + 16, top + 16,
                           art_w - 32, 40, 12.5, False, PALETTE["Muted"])
        tag(note, "ConceptArtifact")
    tag(artifact, "ConceptArtifact")
    if callouts:
        gap = 14
        cw = (art_w - gap * (len(callouts) - 1)) / len(callouts)
        for i, callout in enumerate(callouts):
            label = callout if isinstance(callout, str) else str(callout.get("label", ""))
            role = "highlight" if isinstance(callout, str) else str(callout.get("role", "") or "")
            fill, line = semantic_colors(role, i)
            box = add_card(slide, label, art_left + i * (cw + gap), art_bottom + 14, cw,
                           callout_band - 28, i, 12, fill=fill, line=line)
            tag(box, "ConceptArtifact")


def add_v_model_tool_map(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = set_intro_in_content(slide, s)
    if content is not None:
        tag(content, "VModelToolMap")
    phases = s.get("phases") or s.get("steps") or ["Scope", "Architecture", "Model", "Test", "Deploy"]
    positions = [(78, 185), (238, 276), (410, 354), (582, 276), (742, 185)]
    cards = []
    for i, phase in enumerate(phases[: len(positions)]):
        label = phase if isinstance(phase, str) else str(phase.get("label", ""))
        tool = "" if isinstance(phase, str) else str(phase.get("tool", "") or "")
        x, y = positions[i]
        card = add_card(slide, label, x, y, 136, 54, i, 13.5, fill=PALETTE["LightBlue"], line=PALETTE["Blue"])
        tag(card, "VModelToolMap")
        cards.append((card, x, y))
        if tool:
            tool_box = add_card(slide, tool, x + 12, y + 62, 112, 24, i, 9.5,
                                fill=PALETTE["LightYellow"], line=PALETTE["Orange"])
            tag(tool_box, "VModelToolMap")
    for i in range(len(cards) - 1):
        _, ax, ay = cards[i]
        _, bx, by = cards[i + 1]
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(ax + 136), Pt(ay + 27), Pt(bx), Pt(by + 27))
        line.line.color.rgb = PALETTE["Orange"]
        line.line.width = Pt(1.2)
        tag(line, "VModelToolMap")
    for i, callout in enumerate(s.get("callouts", [])):
        label = callout if isinstance(callout, str) else str(callout.get("label", ""))
        box = add_card(slide, label, 110, 438 + i * 35, 700, 28, i, 11.5,
                       fill=PALETTE["LightGreen"], line=PALETTE["Green"])
        tag(box, "VModelToolMap")


def add_process_state_diagram(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = set_intro_in_content(slide, s)
    if content is not None:
        tag(content, "ProcessStateDiagram")
    steps = list(s.get("steps", []))
    count = max(1, len(steps))
    gap, arrow_w = 20, 22
    band_top = BAND_TOP + (34 if s.get("intro") else 0)
    status_h = 26.0
    card_h = min(190.0, BAND_BOTTOM - band_top - status_h - 6)
    top = band_top + status_h + 6 + (BAND_BOTTOM - band_top - status_h - 6 - card_h) / 2
    width = (BAND_W - (gap + arrow_w) * (count - 1)) / count
    for i, step in enumerate(steps):
        label = step if isinstance(step, str) else str(step.get("label", ""))
        state = "" if isinstance(step, str) else str(step.get("status", "") or "")
        detail = "" if isinstance(step, str) else str(step.get("detail", "") or "")
        fill, line = semantic_colors(state, i)
        x = BAND_LEFT + i * (width + gap + arrow_w)
        text = f"{label}\n{detail}" if detail else label
        card = add_card(slide, text, x, top, width, card_h, i, 14, fill=fill, line=line)
        if detail:
            two_line_card_fonts(card, 14, 11)
        tag(card, "ProcessStateDiagram")
        if state:
            # Status chip sits flush on the card's top edge, not floating above it.
            mark = add_card(slide, state, x, top - status_h + 2, min(96.0, width), status_h, i, 10.5,
                            fill=fill, line=line)
            tag(mark, "ProcessStateDiagram")
        if i < count - 1:
            arrow = add_arrow(slide, x + width + gap / 2, top + card_h / 2 - 10, arrow_w, 20)
            tag(arrow, "ProcessStateDiagram")


def add_demo_exercise(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    panels = s.get("panels") or [
        {"label": "Start", "detail": str(s.get("start", "") or ""), "role": "current"},
        {"label": "Do", "detail": str(s.get("task", "") or ""), "role": "action"},
        {"label": "Keep", "detail": str(s.get("output", "") or ""), "role": "verification"},
    ]
    for i, panel in enumerate(panels):
        label = str(panel.get("label", ""))
        detail = str(panel.get("detail", "") or "")
        role = str(panel.get("role", "") or "")
        if re.search(r"demo|exercise", role):
            fill, line = PALETTE["LightCyan"], PALETTE["Cyan"]
        else:
            fill, line = semantic_colors(role, i)
        text = f"{label}\n{detail}" if detail else label
        gap = 26
        width = (BAND_W - gap * (len(panels) - 1)) / max(1, len(panels))
        longest = max((len(str(p.get("detail", "") or "")) for p in panels), default=0)
        panel_h = min(180.0, max(96.0, 78.0 + longest * 0.55))
        shape = add_card(slide, text, BAND_LEFT + i * (width + gap),
                         BAND_TOP + (BAND_H - panel_h) / 2, width, panel_h, i, 17, fill=fill, line=line)
        if detail:
            two_line_card_fonts(shape, 17, 13)
        tag(shape, "DemoExercise")


def add_recap(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    if s.get("subtitle"):
        # Sits just above the body band, not at a fixed y that leaves an orphan gap.
        box = add_textbox(slide, s["subtitle"], BAND_LEFT, BAND_TOP - 26, BAND_W, 24,
                          18, False, PALETTE["Blue"])
        tag(box, "Recap")
    items = s.get("items") or s.get("bullets") or []
    content = find_content(slide)
    if content is not None:
        top_offset = 40 if s.get("subtitle") else 0
        size = fill_band_with_text(content, items)
        try:
            content.top = Pt(BAND_TOP + top_offset)
            content.height = Pt(BAND_H - top_offset)
        except Exception:
            pass
        set_body_paragraphs(content, items, size)
        space_paragraphs(content, size * 0.7)
        tag(content, "Recap")


def add_section(prs, s, progress_labels):
    """Official Section Header layout: its own title (and text placeholder if present)."""
    slide = new_slide(prs, "Section Header", s)
    outcome = str(s.get("outcome", "") or "")
    if outcome:
        body = find_content(slide)
        if body is not None:
            body.text_frame.text = typographic(outcome)
    return slide


def add_process(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    variant = str(s.get("variant", "horizontal-cards"))
    content = find_content(slide)
    if content is not None:
        set_text(content, str(s.get("intro", "") or ""), 16, False, PALETTE["Muted"])
        tag(content, f"Process_{variant}")
    steps = list(s.get("steps", []))
    count = max(1, len(steps))
    gap = 30
    top = BAND_TOP + (34 if s.get("intro") else 0)
    longest = max((len(str(s2.get("detail", "") or "")) for s2 in steps
                   if isinstance(s2, dict)), default=0)
    height = min(150.0, max(84.0, 66.0 + longest * 0.42))
    top = top + (BAND_BOTTOM - top - height) / 2
    arrow_w = 26
    width = (BAND_W - arrow_w * (count - 1) - gap * (count - 1)) / count
    for i, step in enumerate(steps):
        x = BAND_LEFT + i * (width + gap + arrow_w)
        label = step if isinstance(step, str) else str(step.get("label", ""))
        detail = "" if isinstance(step, str) else str(step.get("detail", "") or "")
        role = "" if isinstance(step, str) else str(step.get("role", "") or "")
        text = f"{label}\n{detail}" if detail else label
        # With a role the fill means something; without one, stay neutral rather than
        # cycling four pastels that encode nothing.
        fill, line = semantic_colors(role, i) if role else (PALETTE["LightBlue"], PALETTE["Blue"])
        card = add_card(slide, text, x, top, width, height, i, 16, fill=fill, line=line)
        if detail:
            two_line_card_fonts(card, 17, 12.5)
        tag(card, f"Process_{variant}")
        if i < count - 1:
            arrow = add_arrow(slide, x + width + gap / 2, top + height / 2 - 11, arrow_w, 22)
            tag(arrow, f"Process_{variant}")


def add_decision(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = find_content(slide)
    if content is not None:
        set_text(content, str(s.get("intro", "") or ""), 16, False, PALETTE["Muted"])
    items = list(s.get("items", []))
    count = max(1, len(items))
    gap = 28
    top = BAND_TOP + (34 if s.get("intro") else 0)
    # A decision card carries what / who proposes / who approves / when / if-no, so it
    # needs room to be the substantial slide the story rules require it to be.
    longest = max((len(str(i if isinstance(i, str) else i.get("label", "")))
                   + len(str("" if isinstance(i, str) else i.get("detail", "") or ""))
                   for i in items), default=0)
    height = min(BAND_BOTTOM - top, max(140.0, 82.0 + longest * 0.36))
    top = top + (BAND_BOTTOM - top - height) / 2
    width = (BAND_W - gap * (count - 1)) / count
    label_size = 18 if longest <= 90 else 15
    for i, item in enumerate(items):
        text = item if isinstance(item, str) else str(item.get("label", ""))
        detail = "" if isinstance(item, str) else str(item.get("detail", "") or "")
        role = "" if isinstance(item, str) else str(item.get("role", "") or "")
        # Without a role the colour cycle can paint an unresolved item green, which
        # states the opposite of the truth under the semantic code. Neutral by default.
        fill, line = semantic_colors(role, i) if role else (PALETTE["LightGray"], PALETTE["Muted"])
        if detail:
            text = f"{text}\n{detail}"
        card = add_card(slide, text, BAND_LEFT + i * (width + gap), top, width, height, i,
                        label_size, fill=fill, line=line)
        if detail:
            two_line_card_fonts(card, label_size, label_size - 4)
        tag(card, "Decision")


def add_artifact_map(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = set_intro_in_content(slide, s)
    if content is not None:
        tag(content, "ArtifactMap")
    artifacts = list(s.get("artifacts", []))
    if not artifacts:
        return
    per_row = grid_shape(len(artifacts))
    row_count = -(-len(artifacts) // per_row)
    top = BAND_TOP + (34 if s.get("intro") else 0)
    rows = band_rows(row_count, 22, top=top, max_height=138)
    gap = 26
    card_w = (BAND_W - gap * (per_row - 1)) / per_row
    for i, artifact in enumerate(artifacts):
        label = artifact if isinstance(artifact, str) else str(artifact.get("label", ""))
        detail = "" if isinstance(artifact, str) else str(artifact.get("detail", "") or "")
        role = "" if isinstance(artifact, str) else str(artifact.get("role", "") or "")
        fill, line = semantic_colors(role, i)
        text = f"{label}\n{detail}" if detail else label
        row, col = divmod(i, per_row)
        row_top, row_h = rows[row]
        # Last row: if it is short, widen its cards to consume the row instead of
        # leaving an orphan hole (a lone card in a 3-wide row is the loudest AI tell).
        in_row = min(per_row, len(artifacts) - row * per_row)
        this_w = (BAND_W - gap * (in_row - 1)) / in_row if in_row < per_row else card_w
        card = add_card(slide, text, BAND_LEFT + col * (this_w + gap), row_top, this_w, row_h, i, 15,
                        fill=fill, line=line)
        if detail:
            two_line_card_fonts(card, 15, 12)
        tag(card, "ArtifactMap")


def add_code_review(prs, s):
    slide = new_slide(prs, "Title and Content", s)
    content = find_content(slide)
    if content is not None:
        set_text(content, str(s.get("intro", "") or ""), 15, False, PALETTE["Muted"])
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Pt(58), Pt(180), Pt(510), Pt(230))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = PALETTE["LightGray"]
    code_box.line.color.rgb = PALETTE["Muted"]
    code_box.shadow.inherit = False
    tf = code_box.text_frame
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    tf.word_wrap = True
    lines = str(s.get("code", "")).split("\n")
    tf.text = lines[0] if lines else ""
    for extra in lines[1:]:
        tf.add_paragraph().text = extra
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = "Menlo"
            run.font.size = Pt(12)
            run.font.color.rgb = PALETTE["Ink"]
    tag(code_box, "CodeReview")
    for i, callout in enumerate(s.get("callouts", [])):
        label = callout if isinstance(callout, str) else str(callout.get("label", ""))
        card = add_card(slide, label, 600, 180 + i * 76, 240, 58, i, 13)
        tag(card, "CodeReview")


# ---------- cover, cleanup, artifacts ----------

def set_cover_text(slide, spec):
    title = None
    subtitle = None
    for shape in slide.shapes:
        ph = _ph_type_name(shape)
        if ph in ("TITLE", "CENTER_TITLE") and title is None:
            title = shape
        elif ph == "SUBTITLE" and subtitle is None:
            subtitle = shape
    if title is None:
        raise SpecError("Title Slide does not expose a title placeholder.")
    title.text_frame.text = str(spec.get("title", ""))
    subtitle_text = str(spec.get("subtitle", "") or "")
    author_line = str(spec.get("authorLine", "") or "")
    date_line = str(spec.get("dateLine", spec.get("date", "")) or "")
    if subtitle_text:
        if subtitle is None:
            raise SpecError("Title Slide does not expose a subtitle placeholder.")
        subtitle.text_frame.text = typographic(subtitle_text)
        # The subtitle inherits a body bullet from the master; a deck subtitle is not a bullet.
        for para in subtitle.text_frame.paragraphs:
            p_pr = para._p.get_or_add_pPr()
            for bullet_tag in ("buChar", "buAutoNum"):
                for node in p_pr.findall(
                        f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{bullet_tag}"):
                    p_pr.remove(node)
            p_pr.append(p_pr.makeelement(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone", {}))
        # Author and date belong in the template's own subtitle placeholder. Drawing
        # them as free-floating boxes would violate the template-precedence rule.
        for line, muted in ((author_line, False), (date_line, True)):
            if not line:
                continue
            para = subtitle.text_frame.add_paragraph()
            para.text = typographic(line)
            p_pr = para._p.get_or_add_pPr()
            p_pr.append(p_pr.makeelement(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone", {}))
            for run in para.runs:
                run.font.size = Pt(12 if muted else 14)
                if muted:
                    run.font.color.rgb = PALETTE["Muted"]
    elif subtitle is not None:
        combined = "\n".join(x for x in (author_line, date_line) if x)
        subtitle.text_frame.text = combined


def delete_slide(prs, slide):
    slide_id_list = prs.slides._sldIdLst
    for sld_id in list(slide_id_list):
        if prs.part.related_part(sld_id.rId) == slide.part:
            prs.part.drop_rel(sld_id.rId)
            slide_id_list.remove(sld_id)
            return


def remove_empty_generated_textboxes(prs):
    for slide in prs.slides:
        for shape in list(slide.shapes):
            try:
                if shape.is_placeholder or not shape.has_text_frame:
                    continue
                if shape.shape_type is not None and shape.shape_type == 17 and not shape.text_frame.text.strip():
                    shape._element.getparent().remove(shape._element)
            except Exception:
                continue


def markdown_list(items):
    lines = []
    for item in items or []:
        if item is None:
            continue
        if isinstance(item, str):
            if item.strip():
                lines.append(f"- {item}")
        elif isinstance(item, dict) and item.get("label"):
            line = f"- {item['label']}"
            if item.get("detail"):
                line += f": {item['detail']}"
            lines.append(line)
        else:
            lines.append(f"- {json.dumps(item)}")
    return "\n".join(lines)


def write_markdown_artifacts(spec, deck_path, template_kind, ground_truth_path, speech_path):
    deck_dir = deck_path.parent
    gt_path = Path(ground_truth_path) if ground_truth_path else deck_dir / "ground-truth.md"
    speech_out = Path(speech_path) if speech_path else deck_dir / "speech.md"
    gt_path.parent.mkdir(parents=True, exist_ok=True)
    speech_out.parent.mkdir(parents=True, exist_ok=True)

    gt = ["# Ground Truth", ""]
    gt.append(f"- Title: {spec.get('title', '')}")
    gt.append(f"- Subtitle: {spec.get('subtitle', '')}")
    gt.append(f"- Use case: {spec.get('useCase', '')}")
    gt.append(f"- Audience: {spec.get('audience', '')}")
    gt.append(f"- Purpose: {spec.get('purpose', '')}")
    gt.append(f"- Template: {template_kind}")
    gt.append(f"- Output: {deck_path}")
    gt.append("- Status: generated from JSON spec; update this ledger when content, assets, or slide order changes.")
    gt += ["", "## Source Facts"]
    gt.append(markdown_list(spec.get("sourceFacts")) or "- Source facts not supplied in spec.")
    gt += ["", "## Progress Labels"]
    gt.append(markdown_list(spec.get("progressLabels")) or "- No shared progress labels supplied.")
    gt += ["", "## Slide Ledger"]
    gt.append("| Slide | Type | Title | Intent | Assets / gaps |")
    gt.append("| --- | --- | --- | --- | --- |")
    gt.append(f"| 1 | cover | {spec.get('title', '')} | Establish title, ownership, and delivery context. | Template cover. |")
    for n, s in enumerate(spec.get("slides", []), start=2):
        gaps = [g for g in (str(s.get("assetStatus", "") or ""), str(s.get("placeholderStatus", "") or ""),
                            markdown_list(s.get("missingInputs")).replace("\n", "<br>")) if g]
        if not gaps:
            gaps = ["No open gap recorded."]
        reason = str(s.get("reason", "")).replace("|", "/")
        title = str(s.get("title", "")).replace("|", "/")
        gt.append(f"| {n} | {s.get('type', '')} | {title} | {reason} | {'<br>'.join(gaps)} |")
    gt += ["", "## Open Inputs"]
    gt.append(markdown_list(spec.get("missingInputs")) or "- None recorded.")

    # The skill requires ground-truth.md to be authored BEFORE the spec exists, so the
    # generated ledger must never clobber the hand-written record. Replace only the
    # generated block; keep everything the author wrote around it.
    generated = "\n".join([GT_BEGIN] + gt + [GT_END])
    if gt_path.exists():
        existing = gt_path.read_text(encoding="utf-8")
        if GT_BEGIN in existing and GT_END in existing:
            head, rest = existing.split(GT_BEGIN, 1)
            _, tail = rest.split(GT_END, 1)
            merged = head + generated + tail
        else:
            merged = existing.rstrip() + "\n\n" + generated + "\n"
        gt_path.write_text(merged, encoding="utf-8")
    else:
        gt_path.write_text(generated + "\n", encoding="utf-8")

    speech = ["# Speech", "", f"## Slide 1 - {spec.get('title', '')}"]
    speech.append(str(spec.get("coverTalkTrack") or
                      "Introduce the topic, audience context, and what the audience should be able to do after the session."))
    for n, s in enumerate(spec.get("slides", []), start=2):
        speech += ["", f"## Slide {n} - {s.get('title', '')}", str(s.get("talkTrack", ""))]
        if s.get("builds"):
            speech += ["", "Build sequence:"]
            speech += [f"- {b}" for b in s["builds"]]
    speech_out.write_text("\n".join(speech), encoding="utf-8")
    return gt_path, speech_out


def export_previews(deck_path, out_dir, scale=2.0):
    soffice = shutil.which("soffice") or "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(deck_path)],
                       check=True, capture_output=True, timeout=300)
        pdf = next(Path(tmp).glob("*.pdf"))
        import pypdfium2 as pdfium
        doc = pdfium.PdfDocument(str(pdf))
        for i, page in enumerate(doc, start=1):
            bitmap = page.render(scale=scale)
            bitmap.to_pil().save(out_dir / f"slide-{i:02d}.png")
        doc.close()
    return out_dir



# ---------- template-first rendering ----------
#
# The point of this skill is to USE the MathWorks template. The template's own
# guidance is explicit: "Avoid manually formatting whenever possible. Instead, use
# built-in styles, templates, layouts, and colors." So every text-bearing preset
# renders through the layout's own title/body placeholders and inherits the
# template's type, sizes, bullets, and colours. Nothing is drawn.

ITEM_FIELDS = ("bullets", "items", "cards", "artifacts", "concepts", "steps",
               "sections", "panels", "phases", "callouts")


def outline_items(s):
    """Collect a slide spec's content into (label, detail) pairs, whatever field it used."""
    out = []
    for field in ITEM_FIELDS:
        for item in s.get(field, []) or []:
            if isinstance(item, str):
                out.append((item, ""))
            else:
                label = str(item.get("label", "") or item.get("tool", "") or "")
                detail = str(item.get("detail", "") or "")
                status = str(item.get("status", "") or item.get("role", "") or "")
                if status and field in ("steps", "panels"):
                    label = f"{label} ({status})" if label else status
                if label or detail:
                    out.append((label, detail))
        if out:
            break
    return out


def add_template_slide(prs, s, layout="Title and Content"):
    """Title + body placeholder, using the template's built-in list styles."""
    slide = new_slide(prs, layout, s)
    body = find_content(slide)
    if body is None:
        return slide
    tf = body.text_frame
    tf.word_wrap = True
    pairs = outline_items(s)
    intro = str(s.get("intro", "") or "")
    first = True

    def para(text, level):
        nonlocal first
        if first:
            p0 = tf.paragraphs[0]
            p0.text = typographic(text)
            p0.level = level
            first = False
            return p0
        p1 = tf.add_paragraph()
        p1.text = typographic(text)
        p1.level = level
        return p1

    if intro:
        para(intro, 0)
    for label, detail in pairs:
        if label:
            para(label, 0)
            if detail:
                para(detail, 1)
        elif detail:
            para(detail, 0)
    if first:
        tf.text = ""
    return slide


BUILDERS = {
    # Text content -> the template's own title/body placeholders.
    "agenda": lambda prs, s, ctx: add_template_slide(prs, s),
    "learning-path": lambda prs, s, ctx: add_template_slide(prs, s),
    "what-to-expect": lambda prs, s, ctx: add_template_slide(prs, s),
    "process": lambda prs, s, ctx: add_template_slide(prs, s),
    "process-build-series": lambda prs, s, ctx: add_template_slide(prs, s),
    "progress-sidebar": lambda prs, s, ctx: add_template_slide(prs, s),
    "decision": lambda prs, s, ctx: add_template_slide(prs, s),
    "content": lambda prs, s, ctx: add_template_slide(prs, s),
    "chapter-objectives": lambda prs, s, ctx: add_template_slide(prs, s),
    "recap": lambda prs, s, ctx: add_template_slide(prs, s),
    "recap-bridge": lambda prs, s, ctx: add_template_slide(prs, s),
    "artifact-map": lambda prs, s, ctx: add_template_slide(prs, s),
    "artifact-review": lambda prs, s, ctx: add_template_slide(prs, s),
    "architecture-layer": lambda prs, s, ctx: add_template_slide(prs, s),
    "team-context": lambda prs, s, ctx: add_template_slide(prs, s),
    "concept-artifact": lambda prs, s, ctx: add_template_slide(prs, s),
    "v-model-tool-map": lambda prs, s, ctx: add_template_slide(prs, s),
    "process-state-diagram": lambda prs, s, ctx: add_template_slide(prs, s),
    "demo-exercise": lambda prs, s, ctx: add_template_slide(prs, s),
    "exercise-demo": lambda prs, s, ctx: add_template_slide(prs, s),
    # Official layouts with their own placeholders.
    "section": lambda prs, s, ctx: add_section(prs, s, ctx["progress_labels"]),
    "chapter-divider": lambda prs, s, ctx: add_section(prs, s, ctx["progress_labels"]),
    "two-content": lambda prs, s, ctx: add_two_content(prs, s),
    # Real content objects (native table, native picture), not decoration.
    "code-review-excerpt": lambda prs, s, ctx: add_code_review(prs, s),
    "code-to-model-review": lambda prs, s, ctx: add_code_review(prs, s),
    "image-evidence": lambda prs, s, ctx: add_image_evidence(prs, s, "ImageEvidence", ctx["spec_dir"]),
    "screenshot-evidence": lambda prs, s, ctx: add_image_evidence(prs, s, "ScreenshotEvidence", ctx["spec_dir"]),
    "screenshot-callout": lambda prs, s, ctx: add_image_evidence(prs, s, "ScreenshotCallout", ctx["spec_dir"]),
    "model-screenshot": lambda prs, s, ctx: add_image_evidence(prs, s, "ModelScreenshot", ctx["spec_dir"]),
    "comparison-table": lambda prs, s, ctx: add_comparison_table(prs, s, "ComparisonTable"),
    "comparison-evidence-table": lambda prs, s, ctx: add_comparison_table(prs, s, "ComparisonEvidenceTable"),
    "results-table": lambda prs, s, ctx: add_comparison_table(prs, s, "ResultsTable"),
}


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("spec_path")
    parser.add_argument("output_path")
    parser.add_argument("--template", choices=["public", "confidential"], default="public")
    parser.add_argument("--ground-truth", default=None)
    parser.add_argument("--speech", default=None)
    parser.add_argument("--export-preview", action="store_true")
    args = parser.parse_args()

    spec_path = Path(args.spec_path).resolve()
    spec_dir = spec_path.parent
    output_path = Path(args.output_path).resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    template_path = SKILL_ROOT / "assets" / "li-mathworks-presentation" / "pptx" / f"{args.template}.pptx"
    if not template_path.exists():
        raise SpecError(f"Template not found: {template_path}")

    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    check_visible_text(spec)
    progress_labels = validate_spec(spec, spec_dir)

    prs = Presentation(str(template_path))
    cover = prs.slides[0]
    if cover.slide_layout.name != "Title Slide":
        raise SpecError("Template first slide must use Title Slide layout.")
    resource_slides = list(prs.slides)[1:]
    set_cover_text(cover, spec)

    ctx = {"progress_labels": progress_labels, "spec_dir": spec_dir}
    for slide_spec in spec.get("slides", []):
        slide_type = str(slide_spec.get("type", ""))
        builder = BUILDERS.get(slide_type)
        if builder is None:
            raise SpecError(f"Unsupported slide type: {slide_type}")
        builder(prs, slide_spec, ctx)
        add_notes(prs.slides[-1], slide_spec)

    for slide in resource_slides:
        delete_slide(prs, slide)
    remove_empty_generated_textboxes(prs)

    prs.save(str(output_path))
    gt_path, speech_path = write_markdown_artifacts(spec, output_path, args.template, args.ground_truth, args.speech)

    result = {
        "OutputPath": str(output_path),
        "SlideCount": len(prs.slides._sldIdLst),
        "TemplateKind": args.template,
        "GroundTruthPath": str(gt_path),
        "SpeechPath": str(speech_path),
    }
    if args.export_preview:
        preview_dir = output_path.parent / f"{output_path.stem}-previews"
        export_previews(output_path, preview_dir)
        result["PreviewDir"] = str(preview_dir)
    print(json.dumps(result, indent=2))


if __name__ == "__main__":
    try:
        main()
    except SpecError as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        sys.exit(1)
