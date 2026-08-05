#!/usr/bin/env python3
"""Render PPTX slides to PNG without PowerPoint or LibreOffice.

Draws the slide master, layout, and slide shapes directly (fills, outlines,
text with inherited sizing, pictures, tables) so generated decks can be
inspected visually on macOS/Linux.

This is an APPROXIMATION for review purposes: it is faithful for geometry,
color, text hierarchy, overlap, and density -- the things slide review needs --
but it is not a substitute for a final check in real PowerPoint (gradients,
effects, exact line breaking, and theme font fallbacks may differ).

Usage:
  .venv/bin/python render_pptx_preview.py DECK.pptx OUTDIR [--scale 1.5]
"""

import argparse
import io
import re
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

FONT_DIRS = ["/System/Library/Fonts/Supplemental", "/Library/Fonts", "/System/Library/Fonts"]
FONT_FILES = {
    ("arial", False): "Arial.ttf",
    ("arial", True): "Arial Bold.ttf",
    ("helvetica", False): "Arial.ttf",
    ("helvetica", True): "Arial Bold.ttf",
    ("menlo", False): "Menlo.ttc",
    ("menlo", True): "Menlo.ttc",
    ("consolas", False): "Menlo.ttc",
    ("consolas", True): "Menlo.ttc",
}
_font_cache = {}

# PowerPoint resolves <a:fld type="slidenum"> at display time; the cached run text
# in the XML is the literal placeholder glyph. Substitute so previews match PowerPoint.
SLIDE_NUMBER_GLYPHS = ("‹#›", "<#>")
_current_slide_number = 0


def resolve_fields(text):
    if not text:
        return text
    for glyph in SLIDE_NUMBER_GLYPHS:
        text = text.replace(glyph, str(_current_slide_number))
    return text


def load_font(name, size_px, bold):
    key = (name or "arial").lower().split(",")[0].strip()
    filename = FONT_FILES.get((key, bold)) or FONT_FILES.get((key, False)) or (
        "Arial Bold.ttf" if bold else "Arial.ttf")
    cache_key = (filename, size_px)
    if cache_key in _font_cache:
        return _font_cache[cache_key]
    for directory in FONT_DIRS:
        path = Path(directory) / filename
        if path.exists():
            try:
                font = ImageFont.truetype(str(path), size_px)
                _font_cache[cache_key] = font
                return font
            except Exception:
                continue
    font = ImageFont.load_default()
    _font_cache[cache_key] = font
    return font


# ---------- theme / color resolution ----------

def theme_colors(prs):
    colors = {}
    try:
        theme_part = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
        root = theme_part._element if hasattr(theme_part, "_element") else None
        if root is None:
            from lxml import etree
            root = etree.fromstring(theme_part.blob)
        scheme = root.find(".//a:clrScheme", NS)
        if scheme is None:
            return colors
        for child in scheme:
            tag = child.tag.split("}")[1]
            srgb = child.find("a:srgbClr", NS)
            sysclr = child.find("a:sysClr", NS)
            if srgb is not None:
                colors[tag] = srgb.get("val")
            elif sysclr is not None:
                colors[tag] = sysclr.get("lastClr", "000000")
    except Exception:
        pass
    # PowerPoint maps lt1/dk1 to bg1/tx1 naming in some places
    colors.setdefault("bg1", colors.get("lt1", "FFFFFF"))
    colors.setdefault("tx1", colors.get("dk1", "000000"))
    colors.setdefault("bg2", colors.get("lt2", "EEEEEE"))
    colors.setdefault("tx2", colors.get("dk2", "333333"))
    return colors


def hex_to_rgb(value):
    if not value or len(value) != 6:
        return None
    try:
        return tuple(int(value[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return None


def apply_lum(rgb, lum_mod, lum_off):
    if rgb is None:
        return None
    r, g, b = [c / 255 for c in rgb]
    if lum_mod is not None:
        r, g, b = r * lum_mod, g * lum_mod, b * lum_mod
    if lum_off is not None:
        r, g, b = r + lum_off, g + lum_off, b + lum_off
    return tuple(max(0, min(255, int(round(c * 255)))) for c in (r, g, b))


def color_from_element(el, theme):
    """Resolve an <a:srgbClr>/<a:schemeClr> child container to RGB."""
    if el is None:
        return None
    srgb = el.find("a:srgbClr", NS)
    if srgb is not None:
        rgb = hex_to_rgb(srgb.get("val"))
        return _apply_mods(srgb, rgb)
    scheme = el.find("a:schemeClr", NS)
    if scheme is not None:
        name = scheme.get("val")
        alias = {"lt1": "lt1", "dk1": "dk1", "bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2"}
        key = alias.get(name, name)
        rgb = hex_to_rgb(theme.get(key) or theme.get(name) or "")
        return _apply_mods(scheme, rgb)
    return None


def _apply_mods(node, rgb):
    if rgb is None:
        return None
    lum_mod = node.find("a:lumMod", NS)
    lum_off = node.find("a:lumOff", NS)
    shade = node.find("a:shade", NS)
    tint = node.find("a:tint", NS)
    lm = int(lum_mod.get("val")) / 100000 if lum_mod is not None else None
    lo = int(lum_off.get("val")) / 100000 if lum_off is not None else None
    rgb = apply_lum(rgb, lm, lo)
    if shade is not None:
        factor = int(shade.get("val")) / 100000
        rgb = tuple(int(c * factor) for c in rgb)
    if tint is not None:
        factor = int(tint.get("val")) / 100000
        rgb = tuple(int(c * factor + 255 * (1 - factor)) for c in rgb)
    return rgb


# ---------- geometry ----------

class Canvas:
    def __init__(self, width_px, height_px, scale):
        self.img = Image.new("RGB", (width_px, height_px), (255, 255, 255))
        self.draw = ImageDraw.Draw(self.img)
        self.scale = scale

    def px(self, emu):
        return int(round(Emu(emu).pt * self.scale)) if emu is not None else 0


def shape_box(shape, canvas):
    try:
        left, top = canvas.px(shape.left), canvas.px(shape.top)
        width, height = canvas.px(shape.width), canvas.px(shape.height)
    except Exception:
        return None
    if width <= 0 or height <= 0:
        return None
    return left, top, left + width, top + height


# ---------- text ----------

def inherited_size(shape, para_el, run_el, master_styles, layout_shape):
    if run_el is not None:
        rpr = run_el.find("a:rPr", NS)
        if rpr is not None and rpr.get("sz"):
            return int(rpr.get("sz")) / 100
    if para_el is not None:
        ppr = para_el.find("a:pPr", NS)
        if ppr is not None:
            defrpr = ppr.find("a:defRPr", NS)
            if defrpr is not None and defrpr.get("sz"):
                return int(defrpr.get("sz")) / 100
    lvl = 0
    if para_el is not None:
        ppr = para_el.find("a:pPr", NS)
        if ppr is not None and ppr.get("lvl"):
            lvl = int(ppr.get("lvl"))
    if layout_shape is not None:
        try:
            lst = layout_shape._element.find(".//a:lstStyle", NS)
            if lst is not None:
                lvl_el = lst.find(f"a:lvl{lvl + 1}pPr", NS)
                if lvl_el is not None:
                    defrpr = lvl_el.find("a:defRPr", NS)
                    if defrpr is not None and defrpr.get("sz"):
                        return int(defrpr.get("sz")) / 100
        except Exception:
            pass
    kind = placeholder_kind(shape)
    style = master_styles.get(kind)
    if style is not None:
        lvl_el = style.find(f"a:lvl{lvl + 1}pPr", NS)
        if lvl_el is not None:
            defrpr = lvl_el.find("a:defRPr", NS)
            if defrpr is not None and defrpr.get("sz"):
                return int(defrpr.get("sz")) / 100
    return 18.0


def placeholder_kind(shape):
    try:
        if not shape.is_placeholder:
            return "other"
        name = shape.placeholder_format.type
        name = name.name if name is not None else ""
    except Exception:
        return "other"
    if name in ("TITLE", "CENTER_TITLE"):
        return "title"
    if name in ("BODY", "OBJECT", "SUBTITLE"):
        return "body"
    return "other"


def master_text_styles(master):
    styles = {}
    try:
        txs = master._element.find("p:txStyles", NS)
        if txs is not None:
            styles["title"] = txs.find("p:titleStyle", NS)
            styles["body"] = txs.find("p:bodyStyle", NS)
            styles["other"] = txs.find("p:otherStyle", NS)
    except Exception:
        pass
    return styles


def wrap_text(draw, text, font, max_width):
    if not text:
        return []
    lines = []
    for raw_line in text.split("\n"):
        words = raw_line.split(" ")
        current = ""
        for word in words:
            trial = f"{current} {word}".strip()
            if draw.textlength(trial, font=font) <= max_width or not current:
                current = trial
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def find_layout_placeholder(shape, layout):
    try:
        if not shape.is_placeholder or layout is None:
            return None
        idx = shape.placeholder_format.idx
        for ph in layout.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
    except Exception:
        return None
    return None


def draw_text_frame(canvas, shape, box, theme, master_styles, layout, default_color=(31, 31, 31)):
    try:
        tf = shape.text_frame
    except Exception:
        return
    left, top, right, bottom = box
    body_pr = shape._element.find(".//a:bodyPr", NS)
    anchor = body_pr.get("anchor") if body_pr is not None else None
    def margin(attr, default_emu):
        if body_pr is not None and body_pr.get(attr):
            return canvas.px(int(body_pr.get(attr)))
        return canvas.px(default_emu)
    ml = margin("lIns", 91440)
    mr = margin("rIns", 91440)
    mt = margin("tIns", 45720)
    mb = margin("bIns", 45720)
    text_left = left + ml
    max_width = max(10, (right - mr) - text_left)
    layout_ph = find_layout_placeholder(shape, layout)
    kind = placeholder_kind(shape)

    rendered = []
    for para in tf.paragraphs:
        para_el = para._p
        runs = para.runs
        if not runs:
            text = para.text
            if not text.strip():
                rendered.append(None)
                continue
            size_pt = inherited_size(shape, para_el, None, master_styles, layout_ph)
            rendered.append([(resolve_fields(text), size_pt, False, default_color, "arial")])
            continue
        pieces = []
        for run in runs:
            run_el = run._r
            size_pt = inherited_size(shape, para_el, run_el, master_styles, layout_ph)
            bold = run.font.bold
            if bold is None:
                bold = kind == "title"
            color = default_color
            try:
                if run.font.color and run.font.color.type is not None:
                    if run.font.color.type == 1 and run.font.color.rgb is not None:
                        color = tuple(run.font.color.rgb)
                    else:
                        rpr = run_el.find("a:rPr", NS)
                        fill = rpr.find("a:solidFill", NS) if rpr is not None else None
                        resolved = color_from_element(fill, theme)
                        if resolved:
                            color = resolved
            except Exception:
                pass
            name = run.font.name or "arial"
            pieces.append((resolve_fields(run.text), size_pt, bool(bold), color, name))
        rendered.append(pieces)

    # measure
    line_data = []
    for pieces in rendered:
        if pieces is None:
            line_data.append((None, 10))
            continue
        text = "".join(p[0] for p in pieces)
        size_pt = max(p[1] for p in pieces)
        bold = pieces[0][2]
        color = pieces[0][3]
        name = pieces[0][4]
        font = load_font(name, max(6, int(round(size_pt * canvas.scale))), bold)
        prefix = ""
        para_index = rendered.index(pieces)
        if kind == "body" and text.strip() and _has_bullet(tf.paragraphs[para_index]._p):
            prefix = "•  "
        for line in wrap_text(canvas.draw, prefix + text, font, max_width):
            line_data.append(((line, font, color), int(round(size_pt * canvas.scale * 1.25))))

    total_height = sum(h for _, h in line_data)
    available = (bottom - mb) - (top + mt)
    if anchor == "ctr":
        y = top + mt + max(0, (available - total_height) // 2)
    elif anchor == "b":
        y = max(top + mt, (bottom - mb) - total_height)
    else:
        y = top + mt
    for payload, height in line_data:
        if payload is not None:
            line, font, color = payload
            canvas.draw.text((text_left, y), line, font=font, fill=color)
        y += height


def _has_bullet(para_el):
    ppr = para_el.find("a:pPr", NS)
    if ppr is None:
        return True
    if ppr.find("a:buNone", NS) is not None:
        return False
    return True


# ---------- shapes ----------

def resolve_fill(shape, theme):
    try:
        sp_pr = shape._element.find(".//a:spPr", NS)
        if sp_pr is None:
            sp_pr = shape._element.find(".//p:spPr", NS)
        if sp_pr is None:
            return None
        if sp_pr.find("a:noFill", NS) is not None:
            return None
        solid = sp_pr.find("a:solidFill", NS)
        if solid is not None:
            return color_from_element(solid, theme)
        grad = sp_pr.find("a:gradFill", NS)
        if grad is not None:
            stop = grad.find(".//a:gs", NS)
            return color_from_element(stop, theme) if stop is not None else None
    except Exception:
        return None
    return None


def resolve_line(shape, theme):
    try:
        ln = shape._element.find(".//a:ln", NS)
        if ln is None:
            return None, 0
        if ln.find("a:noFill", NS) is not None:
            return None, 0
        width_emu = int(ln.get("w")) if ln.get("w") else 12700
        solid = ln.find("a:solidFill", NS)
        color = color_from_element(solid, theme) if solid is not None else None
        return color, width_emu
    except Exception:
        return None, 0


def draw_autoshape(canvas, shape, box, theme):
    left, top, right, bottom = box
    fill = resolve_fill(shape, theme)
    line_color, line_w = resolve_line(shape, theme)
    line_px = max(1, int(round(Emu(line_w).pt * canvas.scale))) if line_color else 0
    geom = ""
    try:
        prst = shape._element.find(".//a:prstGeom", NS)
        geom = prst.get("prst") if prst is not None else ""
    except Exception:
        pass
    if geom in ("roundRect",):
        radius = max(2, int(min(right - left, bottom - top) * 0.09))
        canvas.draw.rounded_rectangle([left, top, right, bottom], radius=radius,
                                      fill=fill, outline=line_color, width=line_px)
    elif geom in ("rightArrow", "downArrow", "leftArrow", "upArrow"):
        color = fill or (215, 136, 36)
        mid_y = (top + bottom) // 2
        third = (bottom - top) // 3
        canvas.draw.polygon([
            (left, mid_y - third // 2), (right - (right - left) // 3, mid_y - third // 2),
            (right - (right - left) // 3, top), (right, mid_y),
            (right - (right - left) // 3, bottom), (right - (right - left) // 3, mid_y + third // 2),
            (left, mid_y + third // 2)], fill=color)
    elif geom in ("line", "straightConnector1", "bentConnector3"):
        if line_color:
            canvas.draw.line([left, top, right, bottom], fill=line_color, width=line_px)
    elif geom in ("ellipse",):
        canvas.draw.ellipse([left, top, right, bottom], fill=fill, outline=line_color, width=line_px)
    else:
        if fill or line_color:
            canvas.draw.rectangle([left, top, right, bottom], fill=fill,
                                  outline=line_color, width=line_px)


def draw_picture(canvas, shape, box):
    left, top, right, bottom = box
    try:
        blob = shape.image.blob
        img = Image.open(io.BytesIO(blob)).convert("RGBA")
        img = img.resize((max(1, right - left), max(1, bottom - top)), Image.LANCZOS)
        canvas.img.paste(img, (left, top), img)
    except Exception:
        canvas.draw.rectangle([left, top, right, bottom], outline=(150, 150, 150), width=1)


def draw_table(canvas, shape, box, theme):
    left, top, right, bottom = box
    table = shape.table
    rows = list(table.rows)
    cols = list(table.columns)
    col_w = [canvas.px(c.width) for c in cols]
    row_h = [canvas.px(r.height) for r in rows]
    total_w, total_h = sum(col_w), sum(row_h)
    if total_w <= 0 or total_h <= 0:
        return
    scale_x = (right - left) / total_w
    scale_y = (bottom - top) / total_h
    y = top
    for r, row in enumerate(rows):
        x = left
        h = int(row_h[r] * scale_y)
        for c in range(len(cols)):
            w = int(col_w[c] * scale_x)
            cell = table.cell(r, c)
            fill = None
            try:
                tc_pr = cell._tc.find(".//a:tcPr", NS)
                if tc_pr is not None:
                    solid = tc_pr.find("a:solidFill", NS)
                    fill = color_from_element(solid, theme) if solid is not None else None
            except Exception:
                pass
            canvas.draw.rectangle([x, y, x + w, y + h], fill=fill, outline=(200, 205, 210), width=1)
            text = cell.text
            if text.strip():
                size_pt = 11.0
                bold = False
                color = (31, 31, 31)
                try:
                    run = cell.text_frame.paragraphs[0].runs[0]
                    if run.font.size:
                        size_pt = run.font.size.pt
                    bold = bool(run.font.bold)
                    if run.font.color and run.font.color.type == 1 and run.font.color.rgb:
                        color = tuple(run.font.color.rgb)
                except Exception:
                    pass
                font = load_font("arial", max(6, int(round(size_pt * canvas.scale))), bold)
                lines = wrap_text(canvas.draw, text, font, max(10, w - 10))
                line_h = int(size_pt * canvas.scale * 1.25)
                anchor = ""
                try:
                    anchor = cell.vertical_anchor.name if cell.vertical_anchor is not None else ""
                except Exception:
                    anchor = ""
                if anchor == "MIDDLE":
                    ty = y + max(4, (h - line_h * len(lines)) // 2)
                elif anchor == "BOTTOM":
                    ty = y + max(4, h - line_h * len(lines) - 4)
                else:
                    ty = y + 4
                for line in lines:
                    canvas.draw.text((x + 5, ty), line, font=font, fill=color)
                    ty += line_h
            x += w
        y += h


def render_shape(canvas, shape, theme, master_styles, layout, skip_empty_placeholders=True):
    box = shape_box(shape, canvas)
    if box is None:
        return
    shape_type = None
    try:
        shape_type = shape.shape_type
    except Exception:
        pass
    if str(shape_type) == "GROUP (6)" or (shape_type is not None and getattr(shape_type, "value", None) == 6):
        for child in shape.shapes:
            render_shape(canvas, child, theme, master_styles, layout, skip_empty_placeholders)
        return
    try:
        if shape.has_table:
            draw_table(canvas, shape, box, theme)
            return
    except Exception:
        pass
    if shape.__class__.__name__ == "Picture":
        draw_picture(canvas, shape, box)
        return
    draw_autoshape(canvas, shape, box, theme)
    try:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if text.strip():
                draw_text_frame(canvas, shape, box, theme, master_styles, layout)
    except Exception:
        pass


def slide_background(canvas, slide, theme):
    for source in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            bg = source._element.find("p:cSld/p:bg", NS)
            if bg is None:
                continue
            fill = bg.find(".//a:solidFill", NS)
            color = color_from_element(fill, theme)
            if color:
                canvas.draw.rectangle([0, 0, canvas.img.width, canvas.img.height], fill=color)
                return
        except Exception:
            continue


def render_deck(deck_path, out_dir, scale=1.5):
    prs = Presentation(str(deck_path))
    theme = theme_colors(prs)
    width_px = int(round(Emu(prs.slide_width).pt * scale))
    height_px = int(round(Emu(prs.slide_height).pt * scale))
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    written = []
    global _current_slide_number
    for index, slide in enumerate(prs.slides, start=1):
        _current_slide_number = index
        canvas = Canvas(width_px, height_px, scale)
        layout = slide.slide_layout
        master = layout.slide_master
        styles = master_text_styles(master)
        slide_background(canvas, slide, theme)
        for shape in master.shapes:
            if shape.is_placeholder:
                continue
            render_shape(canvas, shape, theme, styles, layout)
        for shape in layout.shapes:
            if shape.is_placeholder:
                continue
            render_shape(canvas, shape, theme, styles, layout)
        for shape in slide.shapes:
            render_shape(canvas, shape, theme, styles, layout)
        path = out_dir / f"slide-{index:02d}.png"
        canvas.img.save(path)
        written.append(path)
    return written


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck")
    parser.add_argument("outdir")
    parser.add_argument("--scale", type=float, default=1.5)
    args = parser.parse_args()
    files = render_deck(args.deck, args.outdir, args.scale)
    print(f"Rendered {len(files)} slides to {args.outdir}")


if __name__ == "__main__":
    main()
